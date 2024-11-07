import os

import predibase
from fastapi import  APIRouter, requests,UploadFile, File, HTTPException, Depends, requests
from fastapi.security import HTTPBasicCredentials, HTTPBasic
from predibase import Predibase,FinetuningConfig, DeploymentConfig

from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF
import requests

from myapp import models


# router=APIRouter()

security = HTTPBasic()
users = {
    "IP": "123"
}
def authenticate(credentials: HTTPBasicCredentials = Depends(security)):
    if credentials.username not in users or users[credentials.username] != credentials.password:
        raise HTTPException(status_code=401, detail="Invalid username or password")
    return credentials


#loading the api from predibase
os.environ["PREDIBASE_API_TOKEN"] = "pb_GW8Li822WuK_GutDVN5UBg"

pb = Predibase(api_token="pb_GW8Li822WuK_GutDVN5UBg")
lorax_client = pb.deployments.client("mistral-7b-instruct-v0-2")

# openai.api_key = "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"




# Storage folder
storage_folder = "storage"

# Create storage folder if it doesn't exist
if not os.path.exists(storage_folder):
    os.makedirs(storage_folder)


def extract_text_from_docx(file_path):
    docx_document = Document(file_path)
    text = "\n".join([para.text for para in docx_document.paragraphs])
    return text

def extract_text_from_pptx(file_path):
    pptx_document = Presentation(file_path)
    text = ""
    for slide in pptx_document.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_pdf(file_path):
    pdf_document = fitz.open(file_path)
    text = ""
    for page_num in range(len(pdf_document)):
        page = pdf_document[page_num]
        text += page.get_text()
    return text



def generate_summary(text):
    input_prompt = (
        "[INST] Summarize the following text with a focus on the company's growth prospects, key strategic moves, "
        "and financial highlights for investors. Keep the summary concise and highlight the most impactful information "
        "for next year's earnings: "
        f"{text} [/INST]"
    )
    lorax_client = pb.deployments.client("mistral-7b-instruct-v0-2")
    response = lorax_client.generate(input_prompt)
    summary = response.generated_text
    return summary

# openai.api_key = "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"

# def generate_summary(text):
#     input_prompt = (
#         "Provide a concise summary for an investor, focusing on the company's future growth prospects, "
#         "key changes in the business, potential triggers, and any information that might significantly impact "
#         "next year's earnings and growth. Summarize the following text in 3 sentences: "
#         f"{text}"
#     )
#
#     response = openai.Completion.create(
#         model="text-davinci-003",  # or you can use "gpt-3.5-turbo" for better conversational models
#         prompt=input_prompt,
#         max_tokens=150,
#         temperature=0.5,
#         top_p=1,
#         frequency_penalty=0,
#         presence_penalty=0
#     )
#
#     summary = response.choices[0].text.strip()  # Access the text part of the response
#     return summary

def check_file_extension(file: UploadFile):
    file_extension = file.filename.split(".")[-1]
    if file_extension not in ["docx", "pptx", "pdf"]:
        raise HTTPException(status_code=400, detail="Invalid file type. Only .docx, .pptx, and .pdf files are allowed.")

def check_if_file_already_uploaded(file: UploadFile):
    existing_file = models.Filee.objects.filter(filename=file.filename).first()
    if existing_file:
        raise HTTPException(status_code=400, detail="File already uploaded.")

def save_file_to_storage_folder(file: UploadFile):
    file_path = os.path.join(storage_folder, file.filename)
    with open(file_path, "wb") as f:
        f.write(file.file.read())
    return file_path

def extract_text_from_file(file_path: str, file_extension: str):
    if file_extension == "docx":
        return extract_text_from_docx(file_path)
    elif file_extension == "pptx":
        return extract_text_from_pptx(file_path)
    elif file_extension == "pdf":
        return extract_text_from_pdf(file_path)

def generate_file_summary(text: str):
    summary_response = generate_summary(text)
    return summary_response
















































































# @router.post("/v1/files")
# def upload_file(file: UploadFile = File(...)):
#     # Check the file extension
#     file_extension = file.filename.split(".")[-1]
#     if file_extension not in ["docx", "pptx", "pdf"]:
#         raise HTTPException(status_code=400, detail="Invalid file type. Only.docx,.pptx, and.pdf files are allowed.")
#
#     # Check if file is already uploaded
#     existing_file = models.File.objects.filter(filename=file.filename).first()
#     if existing_file:
#         raise HTTPException(status_code=400, detail="File already uploaded.")
#
#     # Save file to storage folder
#     file_path = os.path.join(storage_folder, file.filename)
#     with open(file_path, "wb") as f:
#         f.write(file.file.read())
#
#
#     if file_extension == "docx":
#         text = extract_text_from_docx(file_path)
#     elif file_extension == "pptx":
#         text = extract_text_from_pptx(file_path)
#     elif file_extension == "pdf":
#         text = extract_text_from_pdf(file_path)
#     # Get file summary from Predibase
#     summary_response = generate_summary(text)
#
#     summary = summary_response
#     new_file=models.File.objects.create(filename=file.filename, filesummary=summary)
#
#     return {"message": "File summary generated successfully"}


# @router.get("/v1/files")
# def list_files(credentials: HTTPBasicCredentials = Depends(security)):
#     # Check username and password
#     if credentials.username not in users or users[credentials.username]!= credentials.password:
#         raise HTTPException(status_code=401, detail="Invalid username or password")
#
#     files = list(models.File.objects.all().values('id', 'filename','filesummary'))  # Use the DjongoManager
#     if not files:
#         return {"message": "No files found"}  # Return an empty message if no files exist
#     return files

# Endpoint 3: Get file summary
# @router.get("/v1/files/{file_id}",response_model=FileSummaryOut)
# def get_file_summary(file_id: str, credentials: HTTPBasicCredentials = Depends(security)):
#     # Check username and password
#     if credentials.username not in users or users[credentials.username] != credentials.password:
#         raise HTTPException(status_code=401, detail="Invalid username or password")
#
#     file = models.File.objects.filter(id=file_id).first()
#
#     if file:
#         return file
#     else:
#         raise HTTPException(status_code=404, detail="File not found")


